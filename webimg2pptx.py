#   Copyright 2023, 2024 hidenorly
#
#   Licensed under the Apache License, Version 2.0 (the "License");
#   you may not use this file except in compliance with the License.
#   You may obtain a copy of the License at
#
#       http://www.apache.org/licenses/LICENSE-2.0
#
#   Unless required by applicable law or agreed to in writing, software
#   distributed under the License is distributed on an "AS IS" BASIS,
#   WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#   See the License for the specific language governing permissions and
#   limitations under the License.

import argparse
import os
import re
import random
import requests
import string
import time

from ImageUtil import ImageUtil

import urllib.request
from urllib.parse import urljoin
from urllib.parse import urlparse
from selenium import webdriver
from selenium.webdriver.common.by import By
from selenium.webdriver.support.ui import WebDriverWait
from selenium.webdriver.support import expected_conditions as EC

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.dml import MSO_THEME_COLOR_INDEX
from pptx.dml.color import RGBColor

import webcolors

globalCache = {}

class UrlUtil:
    def isSameDomain(url1, url2, baseUrl=""):
        isSame = urlparse(url1).netloc == urlparse(url2).netloc
        isbaseUrl =  ( (baseUrl=="") or url2.startswith(baseUrl) )
        return isSame and isbaseUrl

    def getFilenameFromUrl(url):
        filename = ""
        pos = url.find("?")
        if pos!=-1:
            url = url[0:pos]
        pos = url.rfind("/")
        if pos!=-1:
            filename = url[pos+1:]
        return str(filename)


    def get_extension_from_mime(mime_type):
        mime_to_extension = {
            'image/jpeg': 'jpg',
            'image/png': 'png',
            'image/gif': 'gif',
            'image/bmp': 'bmp',
            'image/webp': 'webp',
            'image/svg+xml': 'svg',
            'image/tiff': 'tiff',
            'image/x-icon': 'ico'
        }

        ext = mime_to_extension.get(mime_type, None)
        if ext:
            ext = "."+ext
        return ext

    def getExtFromUrl(url):
        ext=""
        filename = UrlUtil.getFilenameFromUrl(url)
        pos = filename.rfind(".")
        if pos!=-1:
            ext = filename[pos:]

        # fallback if url doesn't contain the file extension
        if not ext:
            try:
                response = requests.head(url)
                content_type = response.headers.get('Content-Type')
                ext = UrlUtil.get_extension_from_mime(content_type)
            except:
                pass

        return str(ext)

    def isValidUrl(url):
        return str(url).startswith("http")


class WebPageImageDownloader:
    def __init__(self, width=1920, height=1080):
        options = webdriver.ChromeOptions()
        options.add_argument('--headless')
        tempDriver = webdriver.Chrome(options=options)
        userAgent = tempDriver.execute_script("return navigator.userAgent")
        userAgent = userAgent.replace("headless", "")
        userAgent = userAgent.replace("Headless", "")

        options = webdriver.ChromeOptions()
        options.add_argument('--headless')
        options.add_argument(f"user-agent={userAgent}")
        driver = webdriver.Chrome(options=options)
        driver.set_window_size(width, height)
        self.driver = driver
        self._driver = tempDriver

    def close(self):
            if self.driver:
                try:
                    self.driver.close()
                except:
                    pass
                self.driver = None
            if self._driver:
                try:
                    self._driver.close()
                except:
                    pass
                self._driver = None

    def getRandomFilename(self):
        letters = string.ascii_lowercase
        return ''.join(random.choice(letters) for i in range(10))

    def getSanitizedFilenameFromUrl(self, url):
        parsed_url = urllib.parse.urlparse(url)
        filename = parsed_url.path.split('/')[-1]

        filename = re.sub(r'[\\/:*?"<>|]', '', filename)

        return filename

    def getOutputFileStream(self, outputPath, url):
        f = None
        filename = self.getSanitizedFilenameFromUrl(url)
        filename = str(os.path.join(outputPath, filename))
        if not filename.endswith(('.png', '.jpg', '.jpeg', '.svg', '.gif', '.webp', '.apng', '.avif')):
            ext =  UrlUtil.getExtFromUrl(url)
            if not ext:
                ext =".jpeg"
            filename = filename+ext

        if os.path.exists(filename):
            fileExt = UrlUtil.getExtFromUrl(filename)
            filename = os.path.join(outputPath, self.getRandomFilename())+fileExt
        try:
            f = open(filename, 'wb')
        except:
            filename = None
            f = None
        filePath = filename
        filename = os.path.basename(filename)
        return f, filename, filePath

    def fallbackDownloadImage(self, imageUrl, outputPath, withFullArgUrl=False):
        filePath = None
        filename = None
        url = None

        try:
            if not withFullArgUrl:
                pos = imageUrl.find("?")
                if pos!=-1:
                    imageUrl = imageUrl[0:pos]
            if UrlUtil.isValidUrl(imageUrl):
                self.driver.get(imageUrl)
                _filename = UrlUtil.getFilenameFromUrl(imageUrl)+".png"
                filePath=os.path.join(outputPath, _filename)
                if os.path.exists(filePath):
                    _filename = self.getRandomFilename()+".png"
                    filePath=os.path.join(outputPath, _filename)
                self.driver.save_screenshot(filePath)
                if os.path.exists(filePath):
                    url = imageUrl
                    filename = _filename

        except Exception as e:
            print(f"Error while processing {imageUrl}: {e}")

        return filename, url, filePath


    def downloadImage(self, imageUrl, outputPath, minDownloadSize=None, withFullArgUrl=False):
        filename = None
        url = None
        if UrlUtil.isValidUrl(imageUrl) and not imageUrl in globalCache:
            globalCache[imageUrl] = True
            filePath = None

            ext = UrlUtil.getExtFromUrl(imageUrl)
            if ext.endswith((".heic", ".HEIC", ".svg", ".webp", ".avif")):
                try:
                    with urllib.request.urlopen(imageUrl) as response:
                        imgContent = response.read()
                        url =imageUrl
                        f, filename, filePath = self.getOutputFileStream(outputPath, imageUrl)
                        if f:
                            f.write(imgContent)
                            f.close()
                except:
                    pass

                if not filePath or not os.path.exists(filePath):
                    # fallback...
                    print(f'Failed to download {imageUrl}')
                    _filename, _url, filePath = self.fallbackDownloadImage(imageUrl, outputPath, withFullArgUrl)
                    if _filename and _url:
                        filename = _filename
                        url = _url

                if filePath and os.path.exists(filePath):
                    if ext.endswith((".svg")):
                        newPngPath = filePath+".png"
                        ImageUtil.convertSvgToPng(filePath, newPngPath)
                        if os.path.exists(newPngPath):
                            filename = newPngPath
                    else:
                        newPath = None
                        if ext.endswith((".webp", ".avif")):
                            # to .png
                            newPath = ImageUtil.covertToPng(filePath)
                        else:
                            # to .jpeg
                            newPath = ImageUtil.covertToJpeg(filePath)
                        if os.path.exists(newPath):
                            size = ImageUtil.getImageSize(newPath)
                            if minDownloadSize==None or (size and size[0] >= minDownloadSize[0] and size[1] >= minDownloadSize[1]):
                                filename = newPath
            else:
                # .png, .jpeg, etc.
                size = None
                response = None
                try:
                    response = requests.get(imageUrl)
                    if response.status_code == 200:
                        # check image size
                        size = ImageUtil.getImageSizeFromChunk(response.content)
                except:
                    print(f'failed to get image size at {imageUrl}')

                if response and response.status_code == 200:
                    if minDownloadSize==None or (size and size[0] >= minDownloadSize[0] and size[1] >= minDownloadSize[1]):
                        url =imageUrl
                        f, filename, filePath = self.getOutputFileStream(outputPath, imageUrl)
                        if f:
                            for chunk in response.iter_content(chunk_size=8192):
                                f.write(chunk)
                            f.close()
                else:
                    # fallback...
                    print(f'Failed to download {imageUrl}')
                    _filename, _url, filePath = self.fallbackDownloadImage(imageUrl, outputPath, withFullArgUrl)
                    if _filename and _url:
                        filename = _filename
                        url = _url

        return filename, url


    def _downloadImagesFromWebPage(self, fileUrls, pageUrls, pageUrl, outputPath, minDownloadSize, baseUrl, maxDepth, depth, usePageUrl, timeOut, withFullArgUrl, scrollPauseTime = 2):
        driver = self.driver
        _imageUrls=[]
        _pageUrls=[]

        if driver==None or depth > maxDepth:
            return

        if not pageUrl in globalCache:
            #globalCache[pageUrl] = True
            try:
                driver.get(pageUrl)
                last_height = driver.execute_script("return document.body.scrollHeight")

                while True:
                    driver.execute_script("window.scrollTo(0, document.body.scrollHeight)")
                    time.sleep(scrollPauseTime)

                    WebDriverWait(driver, timeOut).until(EC.presence_of_element_located((By.TAG_NAME, 'a')))
                    WebDriverWait(driver, timeOut).until(EC.presence_of_element_located((By.TAG_NAME, 'img')))

                    # download image
                    for img_tag in driver.find_elements(By.TAG_NAME, 'img'):
                        imageUrl = None
                        try:
                            imageUrl = img_tag.get_attribute('src')
                        except:
                            pass
                        if imageUrl:
                            imageUrl = urljoin(pageUrl, imageUrl)
                            _imageUrls.append(imageUrl)

                    # get links to other pages
                    links = driver.find_elements(By.TAG_NAME, 'a')
                    for link in links:
                        if link:
                            href = None
                            try:
                                href = link.get_attribute('href')
                            except:
                                continue #print("Error occured (href is not found in a tag) at "+str(link))
                            if href and UrlUtil.isSameDomain(pageUrl, href, baseUrl):
                                if not href in pageUrls:
                                    pageUrls.add(href)
                                    ext = UrlUtil.getExtFromUrl(href)
                                    if ext.endswith(('.png', '.jpg', '.jpeg', '.svg', '.gif', '.webp', '.avif')):
                                        _imageUrls.append(href)
                                    else:
                                        _pageUrls.add(href)
                    new_height = driver.execute_script("return document.body.scrollHeight")
                    if new_height == last_height:
                        break
                    last_height = new_height
            except Exception as e:
                pass #print(f"Error while processing {pageUrl}: {e}")


            for imageUrl in _imageUrls:
                fileName, url = self.downloadImage(imageUrl, outputPath, minDownloadSize, withFullArgUrl)
                if fileName and not fileName in fileUrls:
                    if usePageUrl:
                        fileUrls[fileName] = pageUrl
                    elif url:
                        fileUrls[fileName] = url

            for href in _pageUrls:
                self._downloadImagesFromWebPage(fileUrls, pageUrls, href, outputPath, minDownloadSize, baseUrl, maxDepth, depth + 1, usePageUrl, timeOut, withFullArgUrl)


    def downloadImagesFromWebPages(self, urls, outputPath, minDownloadSize=None, baseUrl="", maxDepth=1, usePageUrl=False, timeOut=60, withFullArgUrl=False):
        fileUrls = {}

        driver = self.driver

        pageUrls=set()
        for url in urls:
            self._downloadImagesFromWebPage(fileUrls, pageUrls, url, outputPath, minDownloadSize, baseUrl, maxDepth, 0, usePageUrl, timeOut, withFullArgUrl)

        return fileUrls


class PowerPointUtil:
    SLIDE_WIDTH_INCH = 16
    SLIDE_HEIGHT_INCH = 9

    def __init__(self, path):
        self.prs = Presentation()
        self.prs.slide_width  = Inches(self.SLIDE_WIDTH_INCH)
        self.prs.slide_height = Inches(self.SLIDE_HEIGHT_INCH)
        self.path = path

    def save(self):
        self.prs.save(self.path)

    # layout is full, left, right, top, bottom
    def getLayoutPosition(self, layout="full"):
        # for full
        x=0
        y=0
        width = self.prs.slide_width
        height = self.prs.slide_height

        if layout=="left" or layout=="right":
            width = width /2
        if layout=="top" or layout=="bottom":
            height = height /2
        if layout=="right":
            x=width
        if layout=="bottom":
            y=height

        return x,y,width,height

    def getLayoutToFitRegion(self, width, height, regionWidth, regionHeight):
        resultWidth = width
        resultHeight = height

        if width > height:
            resultWidth = regionWidth
            resultHeight = int(regionWidth * height / width+0.99)
        else:
            resultHeight = regionHeight
            resultWidth = int(regionHeight * width / height+0.99)

        return resultWidth, regionHeight


    def addSlide(self, layout=None):
        if layout == None:
            layout = self.prs.slide_layouts[6]
        self.currentSlide = self.prs.slides.add_slide(layout)

    def addPicture(self, imagePath, x=0, y=0, width=None, height=None, isFitToSlide=True, regionWidth=None, regionHeight=None, isFitWihthinRegion=False):
        if not regionWidth:
            regionWidth = self.prs.slide_width
        if not regionHeight:
            regionHeight = self.prs.slide_height
        regionWidth = int(regionWidth+0.99)
        regionHeight = int(regionHeight+0.99)
        pic = None
        try:
            pic = self.currentSlide.shapes.add_picture(imagePath, x, y)
        except:
            print(f'failed to add {imagePath}')
        if pic:
            if width and height:
                pic.width = width
                pic.height = height
            else:
                if isFitToSlide:
                    width, height = pic.image.size
                    picWidth = pic.width
                    picHeight = pic.height
                    if width > height:
                        picWidth = regionWidth
                        picHeight = int(regionWidth * height / width + 0.99)
                    else:
                        picHeight = regionHeight
                        picWidth = int(regionHeight * width / height + 0.99)
                    if isFitWihthinRegion:
                        deltaWidth = picWidth - regionWidth
                        deltaHeight = picHeight - regionHeight
                        if deltaWidth>0 or deltaHeight>0:
                            # exceed the region
                            if deltaWidth > deltaHeight:
                                picWidth = regionWidth
                                picHeight = int(regionWidth * height / width + 0.99)
                            else:
                                picHeight = regionHeight
                                picWidth = int(regionHeight * width / height + 0.99)
                    pic.width = picWidth
                    pic.height = picHeight
        return pic

    def nameToRgb(name):
        result = RGBColor(0,0,0)
        try:
            rgb = webcolors.name_to_rgb(name)
            result = RGBColor(rgb.red, rgb.green, rgb.blue)
        except:
            pass
        return result

    def applyExFormat(exFormat, textbox, font, text_frame):
        exFormats = exFormat.split(",")
        for anFormat in exFormats:
            cmdarg = anFormat.split(":")
            cmd = cmdarg[0]
            val = None
            if len(cmdarg)>=2:
                val = cmdarg[1]
            if cmd=="color":
                font.color.rgb = PowerPointUtil.nameToRgb(val)
            elif cmd=="face":
                font.name = val
            elif cmd=="size":
                font.size = Pt(float(val))
            elif cmd=="bold":
                font.bold = True
            elif cmd=="effect":
                # TODO: fix
                shadow = textbox.shadow
                shadow.visible = True
                shadow.shadow_type = 'outer'
                shadow.style = 'outer'
                shadow.blur_radius = Pt(5)
                shadow.distance = Pt(2)
                shadow.angle = 45
                shadow.color = MSO_THEME_COLOR_INDEX.ACCENT_5
                shadow.transparency = 0

    def addText(self, text, x=Inches(0), y=Inches(0), width=None, height=None, fontFace='Calibri', fontSize=Pt(18), isAdjustSize=True, textAlign = PP_ALIGN.LEFT, isVerticalCenter=False, exFormat=None):
        if width==None:
            width=self.prs.slide_width
        if height==None:
            height=self.prs.slide_height
        width = int(width+0.99)
        height = int(height+0.99)

        textbox = self.currentSlide.shapes.add_textbox(x, y, width, height)
        text_frame = textbox.text_frame
        text_frame.text = text
        font = text_frame.paragraphs[0].font
        font.name = fontFace
        font.size = fontSize
        theHeight = textbox.height

        if exFormat:
            PowerPointUtil.applyExFormat(exFormat, textbox, font, text_frame)
        
        if isAdjustSize:
            text_frame.auto_size = True
            textbox.top = y

        if isVerticalCenter:
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        for paragraph in text_frame.paragraphs:
            paragraph.alignment = textAlign



if __name__ == '__main__':
    parser = argparse.ArgumentParser(description='Download images from web pages', formatter_class=argparse.ArgumentDefaultsHelpFormatter)
    parser.add_argument('pages', metavar='PAGE', type=str, nargs='+', help='Web pages to download images from')
    parser.add_argument('-t', '--temp', dest='tempPath', type=str, default='.', help='Temporary path.')
    parser.add_argument("-o", "--output", default='output.pptx', help="Output PowerPoint file path")
    parser.add_argument("-a", "--addUrl", action='store_true', default=False, help='Add URL to the slide')
    parser.add_argument("-p", "--usePageUrl", action='store_true', default=False, help='Use page URL if possible')
    parser.add_argument("-l", "--layout", action='store', default='full', help='Specify layout full or left or right')
    parser.add_argument("-f", "--fullfit", action='store_true', default=False, help='Specify if want to fit within the slide')
    parser.add_argument("-w", "--withFullArgUrl", action='store_true', default=False, help='Specify if want to use full url with ? argument')
    parser.add_argument('--minSize', type=str, help='Minimum size of images to download (format: WIDTHxHEIGHT)')
    parser.add_argument('--maxDepth', type=int, default=1, help='maximum depth of links to follow')
    parser.add_argument('--baseUrl', type=str, default="", help='Specify base url if you want to restrict download under the baseUrl')
    parser.add_argument('--timeOut', type=int, default=60, help='Specify time out [sec] if you want to change the default')
    parser.add_argument('--offsetX', type=float, default=0, help='Specify offset x (Inch. max 16. float)')
    parser.add_argument('--offsetY', type=float, default=0, help='Specify offset y (Inch. max 9. float)')
    parser.add_argument('--fontFace', type=str, default="Calibri", help='Specify font face if necessary')
    parser.add_argument('--fontSize', type=float, default=18.0, help='Specify font size (pt) if necessary')
    parser.add_argument('--title', type=str, default=None, help='Specify title if necessary')
    parser.add_argument('--titleSize', type=float, default=None, help='Specify title size if necessary')
    parser.add_argument('--titleFormat', type=str, default=None, help='Specify title format if necessary e.g. color:white,face:Calibri,size:40,bold')
    args = parser.parse_args()
    if args.usePageUrl:
        args.addUrl = True

    # --- download 
    minDownloadSize = None
    if args.minSize:
        minDownloadSize = tuple(map(int, args.minSize.split('x')))

    if not os.path.exists(args.tempPath):
        os.makedirs(args.tempPath)

    downloader = WebPageImageDownloader()
    fileUrls = downloader.downloadImagesFromWebPages(args.pages, args.tempPath, minDownloadSize, args.baseUrl, args.maxDepth, args.usePageUrl, args.timeOut, args.withFullArgUrl)
    downloader.close()
    downloader = None

    # --- create power point
    prs = PowerPointUtil( args.output )

    # --- sort per page url
    perPageImgFiles={}
    pageUrls = []
    for filename in fileUrls.keys():
        if filename.endswith(('.png', '.jpg', '.jpeg', '.svg', '.gif')):
            pageUrl = fileUrls[filename]
            if pageUrl:
                if not pageUrl in perPageImgFiles:
                    perPageImgFiles[pageUrl] = []
                    pageUrls.append(pageUrl)
                perPageImgFiles[pageUrl].append(filename)

    pageUrls.sort(key=lambda x: (len(x), x))

    # --- add image file to the slide
    x, y, regionWidth, regionHeight = prs.getLayoutPosition(args.layout)
    offsetX = Inches(args.offsetX)
    offsetY = Inches(args.offsetY)
    x = x + offsetX
    y = y + offsetY
    regionWidth = int( regionWidth - offsetX )
    regionHeight = int( regionHeight - offsetY )

    isFitWihthinRegion = args.fullfit

    fontFace = args.fontFace
    fontSize = Pt(args.fontSize)

    textAlign = PP_ALIGN.LEFT
    if args.layout == "right":
        textAlign = PP_ALIGN.RIGHT

    titleSize = args.offsetY*72.0 #Inch to Pt
    if args.titleSize:
        titleSize = Pt(args.titleSize)
    if titleSize<100 or titleSize>400000:
        titleSize = Pt(40) # fail safe

    titleHeight = offsetY
    if titleHeight==0:
        titleHeight = titleSize

    for aPageUrl in pageUrls:
        for filename in perPageImgFiles[aPageUrl]:
            imagePath = os.path.join(args.tempPath, filename)
            if os.path.exists(imagePath):
                prs.addSlide()
                pic = prs.addPicture(imagePath, x, y, None, None, True, regionWidth, regionHeight, isFitWihthinRegion)
                # Add Title
                if args.title:
                    prs.addText(args.title, x, 0, regionWidth, titleHeight, fontFace, titleSize, True, textAlign, True, args.titleFormat)
                # Add filename(URL) at bottom
                if pic and args.addUrl:
                    text = None
                    if not args.usePageUrl:
                        text = filename
                    if filename in fileUrls:
                        text = fileUrls[filename]
                    if text:
                        # TODO: Calc the 0.4
                        prs.addText(text, x, int(y+regionHeight-Inches(0.4)), regionWidth, Inches(0.4), fontFace, fontSize, True, textAlign)

    # --- save the ppt file
    prs.save()
